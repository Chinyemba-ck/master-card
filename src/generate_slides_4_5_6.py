"""
generate_slides_4_5_6.py

Generates PowerPoint slides 4, 5, and 6 for the CatBoost team's
Inclusive Growth Score (IGS) analysis on Winnsboro, Franklin Parish, LA.

Slides:
  4 – IGS Benchmarking (Local): horizontal bar chart of gaps + pillar score table
  5 – Key Findings: 5 mini time-series line charts (2017–2025 divergence)
  6 – Healthcare Access Context: severity bar chart + correlation table

Usage:
    python src/generate_slides_4_5_6.py

Output:
    CatBoost_IGS_Slides_4_5_6.pptx  (in the repo root)
"""

from __future__ import annotations

import io
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.ticker as ticker
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
REPO_ROOT = Path(__file__).resolve().parent.parent
DATA_DIR = REPO_ROOT / "data"
OUT_PATH = REPO_ROOT / "CatBoost_IGS_Slides_4_5_6.pptx"

# ---------------------------------------------------------------------------
# Colour palette
# ---------------------------------------------------------------------------
RED = RGBColor(0xC0, 0x39, 0x2B)        # Winnsboro (problem)
BLUE = RGBColor(0x29, 0x80, 0xB9)       # Archibald (benchmark)
AMBER = RGBColor(0xE6, 0x7E, 0x22)      # threshold / critical
GREEN = RGBColor(0x27, 0xAE, 0x60)      # strength / positive
GRAY = RGBColor(0x7F, 0x8C, 0x8D)       # neutral / national
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x2C, 0x3E, 0x50)       # dark text
LIGHT_BG = RGBColor(0xF5, 0xF6, 0xFA)  # slide background

HEX_RED   = "#C0392B"
HEX_BLUE  = "#2980B9"
HEX_AMBER = "#E67E22"
HEX_GREEN = "#27AE60"
HEX_GRAY  = "#7F8C8D"
HEX_DARK  = "#2C3E50"
HEX_TEAL  = "#16A085"

# Slide size: widescreen 13.33" × 7.5"
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Utility helpers
# ---------------------------------------------------------------------------

def rgb_to_hex(r: int, g: int, b: int) -> str:
    return f"#{r:02X}{g:02X}{b:02X}"


def _add_textbox(slide, left, top, width, height, text, font_size=11,
                 bold=False, color=DARK, align=PP_ALIGN.LEFT,
                 wrap=True, italic=False):
    """Add a simple textbox to a slide."""
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def _set_slide_bg(slide, color: RGBColor):
    """Set uniform slide background fill colour."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _fig_to_image(fig) -> io.BytesIO:
    """Save matplotlib figure to an in-memory PNG BytesIO."""
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    buf.seek(0)
    plt.close(fig)
    return buf


def _add_chart_image(slide, buf, left, top, width, height):
    """Place a chart image (BytesIO) on the slide."""
    slide.shapes.add_picture(buf, left, top, width, height)


def _add_slide_header(slide, title_text: str, subtitle_text: str = "",
                      source_text: str = ""):
    """Add a consistent header bar + optional source attribution."""
    # Top colour bar
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0), SLIDE_W, Inches(0.55)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK
    bar.line.fill.background()

    # Title text in bar
    _add_textbox(slide, Inches(0.18), Inches(0.05), Inches(10), Inches(0.45),
                 title_text, font_size=16, bold=True, color=WHITE)

    if subtitle_text:
        _add_textbox(slide, Inches(10.5), Inches(0.08), Inches(2.65), Inches(0.40),
                     subtitle_text, font_size=9, color=RGBColor(0xBD, 0xC3, 0xC7),
                     align=PP_ALIGN.RIGHT)

    # Source attribution at bottom
    if source_text:
        _add_textbox(slide, Inches(0.18), Inches(7.2), Inches(12), Inches(0.28),
                     source_text, font_size=7, color=GRAY, italic=True)

    # Bottom thin bar
    bbar = slide.shapes.add_shape(
        1, Inches(0), Inches(7.18), SLIDE_W, Inches(0.04)
    )
    bbar.fill.solid()
    bbar.fill.fore_color.rgb = AMBER
    bbar.line.fill.background()


def _add_callout_box(slide, left, top, width, height, text,
                     bg_color: RGBColor = RGBColor(0xFD, 0xF2, 0xE9),
                     border_color: RGBColor = AMBER,
                     font_size=9.5, bold=False):
    """Add a rounded-corner callout box with coloured border."""
    box = slide.shapes.add_shape(
        5,  # ROUNDED_RECTANGLE
        left, top, width, height
    )
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    box.line.color.rgb = border_color
    box.line.width = Pt(1.2)

    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = DARK


# ---------------------------------------------------------------------------
# Load data
# ---------------------------------------------------------------------------

def load_data():
    franklin = pd.read_csv(DATA_DIR / "franklin_parish_indicators.csv")
    richland = pd.read_csv(DATA_DIR / "richland_parish_indicators.csv")
    return franklin, richland


# ===========================================================================
# SLIDE 4  – IGS Benchmarking (Local)
# ===========================================================================

def build_slide_4(prs: Presentation, franklin: pd.DataFrame, richland: pd.DataFrame):
    """Slide 4: IGS Benchmarking — horizontal bar chart + pillar table."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_bg(slide, WHITE)

    _add_slide_header(
        slide,
        "IGS Benchmarking (Local) — Where Winnsboro Falls Behind Archibald (2025)",
        "Slide 4",
        "Data source: data/franklin_parish_indicators.csv · data/richland_parish_indicators.csv  |  CatBoost Team · March 2026"
    )

    # ---- 2025 row for each tract ----------------------------------------
    w2025 = franklin[franklin["Year"] == 2025].iloc[0]
    a2025 = richland[richland["Year"] == 2025].iloc[0]

    # Indicator display names + column keys
    indicator_map = [
        ("Travel Time",          "Travel_Time"),
        ("Min/Women Biz",        "Min_Women_Biz"),
        ("Early Education",      "Early_Education"),
        ("Net Occupancy",        "Net_Occupancy"),
        ("Labor Engagement",     "Labor_Engagement"),
        ("Real Estate",          "Real_Estate"),
        ("Park Land",            "Park_Land"),
        ("Affordable Housing",   "Affordable_Housing"),
        ("Personal Income",      "Personal_Income"),
        ("Gini Coefficient",     "Gini"),
        ("Small Biz Loans",      "Small_Biz_Loans"),
        ("New Businesses",       "New_Businesses"),
        ("Female > Poverty",     "Female_Poverty"),
        ("Health Insurance",     "Health_Insurance"),
        ("Comm. Diversity",      "Comm_Diversity"),
        ("Internet Access",      "Internet_Access"),
    ]

    labels = []
    gaps = []
    w_scores = []
    a_scores = []

    for name, col in indicator_map:
        w_val = w2025.get(col, np.nan)
        a_val = a2025.get(col, np.nan)
        if pd.isna(w_val) or pd.isna(a_val):
            gap = np.nan
        else:
            gap = float(a_val) - float(w_val)
        labels.append(name)
        gaps.append(gap)
        w_scores.append(float(w_val) if not pd.isna(w_val) else 0)
        a_scores.append(float(a_val) if not pd.isna(a_val) else 0)

    # Sort by gap descending (most behind first)
    order = np.argsort(gaps)[::-1]
    labels_s = [labels[i] for i in order]
    gaps_s   = [gaps[i]   for i in order]
    w_s      = [w_scores[i] for i in order]
    a_s      = [a_scores[i] for i in order]

    bar_colors = [HEX_AMBER if g > 0 else HEX_TEAL for g in gaps_s]

    # ---- Build horizontal bar chart (left panel) -------------------------
    fig, ax = plt.subplots(figsize=(6.0, 5.8), facecolor="white")

    y_pos = np.arange(len(labels_s))
    bars = ax.barh(y_pos, gaps_s, color=bar_colors, height=0.65,
                   edgecolor="white", linewidth=0.4)

    # Value labels
    for i, (g, bar) in enumerate(zip(gaps_s, bars)):
        ha = "left" if g >= 0 else "right"
        offset = 0.8 if g >= 0 else -0.8
        ax.text(g + offset, i, f"{g:+.0f}", va="center", ha=ha,
                fontsize=7.5, color=HEX_DARK, fontweight="bold")

    ax.set_yticks(y_pos)
    ax.set_yticklabels(labels_s, fontsize=8.5)
    ax.set_xlabel("Gap (Archibald − Winnsboro)", fontsize=8.5, color=HEX_DARK)
    ax.set_title("Indicator Gap: Archibald vs Winnsboro (2025)", fontsize=10,
                 fontweight="bold", color=HEX_DARK, pad=8)
    ax.axvline(0, color=HEX_DARK, linewidth=0.8)
    ax.set_xlim(-30, 85)
    ax.tick_params(axis="x", labelsize=7.5)
    ax.spines[["top", "right"]].set_visible(False)
    ax.grid(axis="x", linestyle="--", alpha=0.35)

    # Legend
    patch_behind = mpatches.Patch(color=HEX_AMBER, label="Archibald leads (gap)")
    patch_ahead  = mpatches.Patch(color=HEX_TEAL,  label="Winnsboro leads")
    ax.legend(handles=[patch_behind, patch_ahead], fontsize=7.5,
              loc="lower right", framealpha=0.85)

    fig.tight_layout()
    chart_buf = _fig_to_image(fig)
    _add_chart_image(slide, chart_buf,
                     Inches(0.25), Inches(0.65), Inches(6.6), Inches(6.25))

    # ---- Right panel: pillar score table ---------------------------------
    # Section title
    _add_textbox(slide, Inches(7.1), Inches(0.65), Inches(5.8), Inches(0.32),
                 "Pillar Score Comparison — 2025", font_size=10.5, bold=True,
                 color=DARK)

    # Build pillar table image
    fig2, ax2 = plt.subplots(figsize=(4.5, 3.6), facecolor="white")
    ax2.axis("off")

    pillars = ["Overall IGS", "Place", "Economy", "Community"]
    w_vals  = [
        float(w2025["IGS"]),
        float(w2025["Place"]),
        float(w2025["Economy"]),
        float(w2025["Community"]),
    ]
    a_vals  = [
        float(a2025["IGS"]),
        float(a2025["Place"]),
        float(a2025["Economy"]),
        float(a2025["Community"]),
    ]

    col_labels = ["Pillar", "Winnsboro", "Archibald", "Gap"]
    table_data = []
    for p, wv, av in zip(pillars, w_vals, a_vals):
        table_data.append([p, f"{wv:.0f}", f"{av:.0f}", f"{av-wv:+.0f}"])

    tbl = ax2.table(
        cellText=table_data,
        colLabels=col_labels,
        cellLoc="center",
        loc="center",
        bbox=[0, 0, 1, 1],
    )
    tbl.auto_set_font_size(False)
    tbl.set_fontsize(11)
    for (row, col), cell in tbl.get_celld().items():
        if row == 0:
            cell.set_facecolor(HEX_DARK)
            cell.set_text_props(color="white", fontweight="bold")
        else:
            pillar_name = table_data[row - 1][0]
            if col == 1:   # Winnsboro
                cell.set_facecolor("#FADBD8")
                cell.set_text_props(color=HEX_RED, fontweight="bold")
            elif col == 2: # Archibald
                cell.set_facecolor("#D6EAF8")
                cell.set_text_props(color=HEX_BLUE, fontweight="bold")
            elif col == 3: # Gap
                gap_val = float(table_data[row - 1][3])
                if gap_val < -10:
                    cell.set_facecolor("#FDEDEC")
                    cell.set_text_props(color=HEX_RED, fontweight="bold")
                else:
                    cell.set_facecolor("#FDFEFE")
            cell.set_edgecolor("#D5D8DC")
        cell.set_height(0.22)
    fig2.tight_layout()
    tbl_buf = _fig_to_image(fig2)
    _add_chart_image(slide, tbl_buf,
                     Inches(7.1), Inches(1.0), Inches(5.8), Inches(2.8))

    # ---- Critical gaps callout -------------------------------------------
    _add_callout_box(
        slide,
        Inches(7.1), Inches(3.95), Inches(5.8), Inches(3.1),
        "⚠  CRITICAL GAPS (2025)\n\n"
        "• Travel Time:         9  vs  70  (−61)\n"
        "• Min/Women Biz:     15  vs  82  (−67)\n"
        "• Early Education:   19  vs  69  (−50)\n"
        "• Net Occupancy:     38  vs  85  (−47)\n"
        "• Labor Engagement: 14  vs  58  (−44)\n\n"
        "Winnsboro's collapse is NOT structural — Archibald\n"
        "recovered +21 IGS points in 8 years. These gaps\n"
        "are reversible.",
        bg_color=RGBColor(0xFD, 0xF2, 0xE9),
        border_color=AMBER,
        font_size=9.0,
    )


# ===========================================================================
# SLIDE 5  – Key Findings: Five Indicators That Diverged (2017–2025)
# ===========================================================================

def build_slide_5(prs: Presentation, franklin: pd.DataFrame, richland: pd.DataFrame):
    """Slide 5: Key Findings — 5 mini time-series line charts."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_bg(slide, WHITE)

    _add_slide_header(
        slide,
        "Key Findings: How Winnsboro and Archibald Diverged (2017–2025)",
        "Slide 5",
        "Data source: data/franklin_parish_indicators.csv · data/richland_parish_indicators.csv  |  CatBoost Team · March 2026"
    )

    years = sorted(franklin["Year"].unique())

    indicators_5 = [
        ("Personal Income",     "Personal_Income"),
        ("Net Occupancy",       "Net_Occupancy"),
        ("Labor Engagement",    "Labor_Engagement"),
        ("Small Biz Loans",     "Small_Biz_Loans"),
        ("Early Education",     "Early_Education"),
    ]

    # ---- 5-panel mini line chart -----------------------------------------
    fig, axes = plt.subplots(1, 5, figsize=(13.0, 3.25), facecolor="white",
                             sharey=False)
    fig.subplots_adjust(wspace=0.38)

    for ax, (title, col) in zip(axes, indicators_5):
        w_vals = [float(franklin[franklin["Year"] == y][col].values[0])
                  if not pd.isna(franklin[franklin["Year"] == y][col].values[0])
                  else np.nan for y in years]
        a_vals = [float(richland[richland["Year"] == y][col].values[0])
                  if not pd.isna(richland[richland["Year"] == y][col].values[0])
                  else np.nan for y in years]

        ax.plot(years, a_vals, color=HEX_BLUE, linewidth=2.2,
                marker="o", markersize=3.5, label="Archibald")
        ax.plot(years, w_vals, color=HEX_RED,  linewidth=2.2,
                marker="s", markersize=3.5, label="Winnsboro", linestyle="--")

        ax.axhline(45, color=HEX_AMBER, linewidth=0.9, linestyle=":",
                   label="IGS 45 base")

        # Annotate start and end
        for vals, clr in ((a_vals, HEX_BLUE), (w_vals, HEX_RED)):
            start_v = next((v for v in vals if not np.isnan(v)), None)
            end_v   = next((v for v in reversed(vals) if not np.isnan(v)), None)
            start_y = years[next((i for i, v in enumerate(vals)
                                  if not np.isnan(v)), 0)]
            end_y   = years[next((i for i, v in reversed(list(enumerate(vals)))
                                  if not np.isnan(v)), -1)]
            if start_v is not None:
                ax.annotate(f"{start_v:.0f}", (start_y, start_v),
                            textcoords="offset points", xytext=(-4, 5),
                            fontsize=6.5, color=clr, fontweight="bold")
            if end_v is not None:
                ax.annotate(f"{end_v:.0f}", (end_y, end_v),
                            textcoords="offset points", xytext=(2, 5),
                            fontsize=6.5, color=clr, fontweight="bold")

        ax.set_title(title, fontsize=8.5, fontweight="bold", color=HEX_DARK, pad=5)
        ax.set_xticks(years[::2])
        ax.set_xticklabels([str(y) for y in years[::2]], rotation=45,
                           fontsize=6.5)
        ax.tick_params(axis="y", labelsize=7)
        ax.set_ylim(0, 105)
        ax.spines[["top", "right"]].set_visible(False)
        ax.grid(axis="y", linestyle="--", alpha=0.3)
        ax.set_ylabel("Score (0–100)", fontsize=7)

    # Shared legend
    handles, lbls = axes[0].get_legend_handles_labels()
    fig.legend(handles, lbls, loc="upper center", ncol=3,
               fontsize=8, framealpha=0.85,
               bbox_to_anchor=(0.5, 1.02))
    fig.suptitle("", y=1.05)
    chart_buf = _fig_to_image(fig)
    _add_chart_image(slide, chart_buf,
                     Inches(0.15), Inches(0.65), Inches(13.0), Inches(3.5))

    # ---- Insight callout boxes (bottom row) ------------------------------
    insights = [
        ("Root Trigger",
         "Early Education collapsed\n78 → 19 (−59 pts)\nRoot cause of all downstream failures"),
        ("Primary IGS Lever",
         "Labor Engagement dropped\n48 → 14 (−34 pts)\nSingle largest IGS driver in all 3 models"),
        ("Recovery Proof",
         "Archibald recovered +21 IGS\nin 8 years — these gaps\nare fully reversible"),
        ("Model Confirmation",
         "Ridge R²=0.935 confirms:\nEarly Ed (+0.99) →\nLabor Eng. (+1.86) →\nPersonal Income (+2.15)"),
    ]

    box_w = Inches(3.1)
    box_h = Inches(2.5)
    x_starts = [Inches(0.15), Inches(3.38), Inches(6.61), Inches(9.84)]
    bg_colors = [
        RGBColor(0xFD, 0xF2, 0xE9),
        RGBColor(0xFD, 0xF2, 0xE9),
        RGBColor(0xEB, 0xF5, 0xFB),
        RGBColor(0xEB, 0xF5, 0xFB),
    ]
    border_colors = [AMBER, AMBER, BLUE, BLUE]

    for (heading, body), x, bg, bdr in zip(insights, x_starts,
                                            bg_colors, border_colors):
        box = slide.shapes.add_shape(
            5, x, Inches(4.28), box_w, box_h
        )
        box.fill.solid()
        box.fill.fore_color.rgb = bg
        box.line.color.rgb = bdr
        box.line.width = Pt(1.2)

        tf = box.text_frame
        tf.word_wrap = True

        p0 = tf.paragraphs[0]
        r0 = p0.add_run()
        r0.text = heading
        r0.font.size = Pt(9.5)
        r0.font.bold = True
        r0.font.color.rgb = bdr

        p1 = tf.add_paragraph()
        r1 = p1.add_run()
        r1.text = "\n" + body
        r1.font.size = Pt(8.5)
        r1.font.color.rgb = DARK


# ===========================================================================
# SLIDE 6  – Healthcare Access Context
# ===========================================================================

def build_slide_6(prs: Presentation, franklin: pd.DataFrame, richland: pd.DataFrame):
    """Slide 6: Healthcare Access Context."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_bg(slide, WHITE)

    _add_slide_header(
        slide,
        "Healthcare Access Context: Why Economic Conditions Matter More Than Coverage",
        "Slide 6",
        "Data source: data/franklin_parish_indicators.csv · data/richland_parish_indicators.csv · README.md  |  CatBoost Team · March 2026"
    )

    # ---- Top callout banner ----------------------------------------------
    banner = slide.shapes.add_shape(
        1, Inches(0.15), Inches(0.62), Inches(13.0), Inches(0.52)
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = RGBColor(0xEB, 0xF5, 0xFB)
    banner.line.color.rgb = BLUE
    banner.line.width = Pt(1.2)
    _add_textbox(
        slide, Inches(0.3), Inches(0.64), Inches(12.7), Inches(0.46),
        "93.5% of Winnsboro residents are insured — above the 92.5% national base.  "
        "The problem is NOT coverage.  It's affordability & mobility.",
        font_size=10.5, bold=True, color=BLUE, align=PP_ALIGN.CENTER
    )

    # ---- Left panel: Correlation table -----------------------------------
    w2025 = franklin[franklin["Year"] == 2025].iloc[0]
    a2025 = richland[richland["Year"] == 2025].iloc[0]

    table_rows = [
        # (Indicator, W score, Natl Coeff, Arch r, Gap, Healthcare Impact)
        ("Labor\nEngagement",   "14",   "+1.86", "+0.91", "−44",
         "CRITICAL — no job = no income\nfor copays, transport"),
        ("Early\nEducation",    "19",   "+0.99", "+0.78", "−50",
         "CRITICAL BOTTLENECK — no\nchildcare = parents can't work"),
        ("Personal\nIncome",    "36",   "+2.15", "+0.94", "−20",
         "HIGH — $44K vs $75K national;\ncan't afford deductibles"),
        ("Net\nOccupancy",      "38",   "+1.88", "+0.83", "−47",
         "CRITICAL — families leave\nwhen jobs disappear"),
        ("Internet\nAccess",    "2",    "—",     "—",     "—",
         "CRITICAL → BEING RESOLVED\n(Volt Broadband Oct 2024)"),
    ]

    col_headers = ["Indicator", "W Score", "Natl\nCoeff", "Arch\nr", "Gap",
                   "Healthcare Impact"]

    fig3, ax3 = plt.subplots(figsize=(7.2, 3.9), facecolor="white")
    ax3.axis("off")

    all_data = [col_headers] + table_rows
    # Separate headers from data
    tbl3 = ax3.table(
        cellText=table_rows,
        colLabels=col_headers,
        cellLoc="center",
        loc="center",
        bbox=[0, 0, 1, 1],
    )
    tbl3.auto_set_font_size(False)
    tbl3.set_fontsize(8.5)

    row_colors_bg = [
        "#FADBD8",   # CRITICAL red tint
        "#FADBD8",
        "#FEF9E7",   # HIGH amber tint
        "#FADBD8",
        "#FDFEFE",   # neutral
    ]

    for (row, col), cell in tbl3.get_celld().items():
        cell.set_edgecolor("#D5D8DC")
        if row == 0:
            cell.set_facecolor(HEX_DARK)
            cell.set_text_props(color="white", fontweight="bold")
            cell.set_height(0.20)
        else:
            cell.set_facecolor(row_colors_bg[row - 1])
            if col == 4 and row <= 4:   # Gap column — critical items
                cell.set_text_props(color=HEX_RED, fontweight="bold")
            cell.set_height(0.20)

    # Set column widths
    col_widths = [0.12, 0.10, 0.10, 0.09, 0.08, 0.51]
    for col_idx, cw in enumerate(col_widths):
        for row_idx in range(len(table_rows) + 1):
            tbl3[(row_idx, col_idx)].set_width(cw)

    fig3.tight_layout()
    tbl3_buf = _fig_to_image(fig3)
    _add_chart_image(slide, tbl3_buf,
                     Inches(0.15), Inches(1.22), Inches(7.4), Inches(3.55))

    # ---- Right panel: severity bar chart --------------------------------
    indicators_bar = [
        # (label, W score, category)
        ("Labor Engagement",  14,  "CRITICAL"),
        ("Net Occupancy",     38,  "CRITICAL"),
        ("Travel Time",        9,  "CRITICAL"),
        ("Early Education",   19,  "CRITICAL"),
        ("Personal Income",   36,  "HIGH"),
        ("Internet Access",    2,  "HIGH"),
        ("Health Insurance",  59,  "STRENGTH"),
        ("Female > Poverty",  69,  "STRENGTH"),
        ("New Businesses",    73,  "STRENGTH"),
    ]

    sev_color_map = {
        "CRITICAL":  HEX_RED,
        "HIGH":      HEX_AMBER,
        "STRENGTH":  HEX_GREEN,
    }

    labels_b = [r[0] for r in indicators_bar]
    scores_b = [r[1] for r in indicators_bar]
    colors_b = [sev_color_map[r[2]] for r in indicators_bar]

    fig4, ax4 = plt.subplots(figsize=(5.0, 4.0), facecolor="white")
    y_pos = np.arange(len(labels_b))[::-1]
    ax4.barh(y_pos, scores_b, color=colors_b, height=0.62,
             edgecolor="white", linewidth=0.4)

    for i, (score, yp) in enumerate(zip(scores_b, y_pos)):
        ax4.text(score + 1, yp, f"{score}", va="center", fontsize=8,
                 color=HEX_DARK, fontweight="bold")

    ax4.axvline(45, color=HEX_AMBER, linewidth=1.1, linestyle=":",
                label="IGS 45 threshold")
    ax4.set_yticks(y_pos)
    ax4.set_yticklabels(labels_b[::-1], fontsize=8.5)
    ax4.set_xlabel("Score (0–100)", fontsize=8.5, color=HEX_DARK)
    ax4.set_title("Winnsboro 2025 Indicator Scores\n(Severity-coded)",
                  fontsize=9.5, fontweight="bold", color=HEX_DARK, pad=6)
    ax4.set_xlim(0, 105)
    ax4.spines[["top", "right"]].set_visible(False)
    ax4.grid(axis="x", linestyle="--", alpha=0.3)
    ax4.tick_params(axis="x", labelsize=7.5)

    # Legend
    p_crit = mpatches.Patch(color=HEX_RED,   label="CRITICAL (below 20)")
    p_high = mpatches.Patch(color=HEX_AMBER, label="HIGH (20–45)")
    p_str  = mpatches.Patch(color=HEX_GREEN, label="STRENGTH (above 55)")
    ax4.legend(handles=[p_crit, p_high, p_str], fontsize=7.5,
               loc="lower right", framealpha=0.85)

    fig4.tight_layout()
    bar_buf = _fig_to_image(fig4)
    _add_chart_image(slide, bar_buf,
                     Inches(7.7), Inches(1.22), Inches(5.4), Inches(4.0))

    # ---- Bottom: Archibald Proof callout ---------------------------------
    _add_callout_box(
        slide,
        Inches(0.15), Inches(4.9), Inches(12.95), Inches(2.1),
        "THE ARCHIBALD PROOF  |  Archibald recovered these same indicators, "
        "proving the chain is reversible:\n"
        "  Childcare re-opened  →  Women returned to workforce  →  "
        "Household income rose  →  Families stayed  →  "
        "Housing occupancy stabilised  →  Local economy recovered.\n"
        "  Early Education: 40 → 69  (+29)  |  "
        "Labor Engagement: 38 → 58  (+20)  |  "
        "Net Occupancy: 48 → 85  (+37)  |  "
        "Personal Income: 37 → 56  (+19)  |  "
        "Overall IGS: 48 → 59  (+21)",
        bg_color=RGBColor(0xEA, 0xF7, 0xEE),
        border_color=GREEN,
        font_size=9.0,
    )


# ===========================================================================
# Main
# ===========================================================================

def main():
    print("Loading data …")
    franklin, richland = load_data()

    print("Creating presentation …")
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print("  Building Slide 4: IGS Benchmarking (Local) …")
    build_slide_4(prs, franklin, richland)

    print("  Building Slide 5: Key Findings — Divergence (2017–2025) …")
    build_slide_5(prs, franklin, richland)

    print("  Building Slide 6: Healthcare Access Context …")
    build_slide_6(prs, franklin, richland)

    prs.save(OUT_PATH)
    print(f"\n✓  Saved → {OUT_PATH}")


if __name__ == "__main__":
    main()
